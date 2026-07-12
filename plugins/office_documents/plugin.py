import csv
import io
import os

from agno.tools.toolkit import Toolkit


def _ensure_parent_dir(path: str) -> None:
    """Create the parent directory of path if it doesn't exist yet, mirroring
    pc_control.PCControlTools.write_file's convention."""
    parent = os.path.dirname(os.path.abspath(path))
    if parent:
        os.makedirs(parent, exist_ok=True)


class OfficeDocumentTools(Toolkit):
    """Real Office file generation: Excel, Word, PowerPoint, and PDF. Every
    function always requires confirmation (enforced in plugin_registry.py)."""

    def __init__(self, **kwargs):
        tools = [
            self.create_excel,
            self.create_word_document,
            self.create_powerpoint,
            self.create_pdf,
        ]
        super().__init__(
            name="office_documents",
            tools=tools,
            requires_confirmation_tools=[t.__name__ for t in tools],
            **kwargs,
        )

    def create_excel(self, path: str, csv_data: str, sheet_name: str = "Sheet1") -> str:
        """Create an Excel (.xlsx) file. csv_data is CSV-formatted text (rows separated
        by newlines, columns by commas) - the first row is typically the header."""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = sheet_name
        for row in csv.reader(io.StringIO(csv_data)):
            ws.append(row)
        _ensure_parent_dir(path)
        wb.save(path)
        return f"Created Excel file at {path}"

    def create_word_document(self, path: str, content: str, title: str = "") -> str:
        """Create a Word (.docx) document. content is plain text with paragraphs
        separated by blank lines. title is an optional heading."""
        from docx import Document

        doc = Document()
        if title:
            doc.add_heading(title, level=1)
        for paragraph in content.split("\n\n"):
            if paragraph.strip():
                doc.add_paragraph(paragraph.strip())
        _ensure_parent_dir(path)
        doc.save(path)
        return f"Created Word document at {path}"

    def create_powerpoint(self, path: str, slides_text: str) -> str:
        """Create a PowerPoint (.pptx) file. slides_text holds multiple slides
        separated by a line with just '---'; within a slide, the first line is
        the title and the rest is body text (one bullet per line)."""
        from pptx import Presentation

        prs = Presentation()
        layout = prs.slide_layouts[1]  # title + content
        for slide_block in slides_text.split("---"):
            lines = [line for line in slide_block.strip().split("\n") if line.strip()]
            if not lines:
                continue
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = lines[0]
            if len(lines) > 1:
                body = slide.placeholders[1].text_frame
                body.text = lines[1]
                for line in lines[2:]:
                    p = body.add_paragraph()
                    p.text = line
        _ensure_parent_dir(path)
        prs.save(path)
        return f"Created PowerPoint file at {path}"

    def create_pdf(self, path: str, content: str, title: str = "") -> str:
        """Create a PDF file. content is plain text with paragraphs separated by
        blank lines. title is an optional heading at the top."""
        from fpdf import FPDF

        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Helvetica", size=12)
        if title:
            pdf.set_font("Helvetica", "B", 16)
            pdf.multi_cell(0, 10, title)
            pdf.ln(4)
            pdf.set_font("Helvetica", size=12)
        for paragraph in content.split("\n\n"):
            if paragraph.strip():
                pdf.multi_cell(0, 8, paragraph.strip())
                pdf.ln(4)
        _ensure_parent_dir(path)
        pdf.output(path)
        return f"Created PDF file at {path}"


def get_toolkit():
    return OfficeDocumentTools()
